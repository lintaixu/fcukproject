"""
Create professional PPTX for K-line pattern analysis project
Based on fcu.pptx.pdf structure
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
import copy

# ── Color Palette (Midnight Executive) ────────────────────────────────────────
NAVY    = RGBColor(0x1E, 0x27, 0x61)   # primary dark
ICE     = RGBColor(0xCA, 0xDC, 0xFC)   # light blue
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_L  = RGBColor(0xF4, 0xF6, 0xFB)   # light background
GRAY_M  = RGBColor(0x94, 0xA3, 0xB8)   # muted text
ACCENT  = RGBColor(0x4A, 0x90, 0xE2)   # accent blue
RED     = RGBColor(0xE5, 0x3E, 0x3E)
GREEN   = RGBColor(0x38, 0xA1, 0x69)
GOLD    = RGBColor(0xD4, 0xAF, 0x37)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]  # blank

# ── Helper Functions ───────────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill_color, line_color=None, line_width=0):
    from pptx.util import Pt
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h, font_size=14, color=WHITE, bold=False,
             align=PP_ALIGN.LEFT, font_face="Calibri", italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_face
    run.font.italic = italic
    return txBox

def add_multiline(slide, lines, x, y, w, h, font_size=14, color=WHITE,
                  bold=False, align=PP_ALIGN.LEFT, spacing=1.15):
    """lines: list of (text, bold, size, color) or just strings"""
    from pptx.util import Pt
    from pptx.oxml.ns import qn
    from lxml import etree
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if isinstance(line, str):
            txt, b, sz, c = line, bold, font_size, color
        else:
            txt, b, sz, c = line
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = txt
        run.font.size = Pt(sz)
        run.font.color.rgb = c
        run.font.bold = b
        run.font.name = "Calibri"
    return txBox

def nav_bar(slide, current_num, total=9):
    """thin bottom nav bar with slide number"""
    add_rect(slide, 0, 7.2, 13.33, 0.3, NAVY)
    add_text(slide, f"{current_num} / {total}", 12.3, 7.2, 1.0, 0.3,
             font_size=10, color=GRAY_M, align=PP_ALIGN.RIGHT)

def section_header(slide, text, subtitle=""):
    """left-side vertical accent bar + title"""
    add_rect(slide, 0.5, 0.3, 0.08, 0.55, ACCENT)
    add_text(slide, text, 0.7, 0.28, 11, 0.6,
             font_size=26, color=NAVY, bold=True, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle, 0.7, 0.85, 11, 0.35,
                 font_size=13, color=GRAY_M, align=PP_ALIGN.LEFT)

def card(slide, x, y, w, h):
    """white card with subtle shadow effect"""
    # shadow
    add_rect(slide, x+0.04, y+0.04, w, h, RGBColor(0xD0, 0xD8, 0xE8))
    add_rect(slide, x, y, w, h, WHITE)

# ─────────────────────────────────────────────────────────────────────────────
# Slide 1: Title
# ─────────────────────────────────────────────────────────────────────────────
sl1 = prs.slides.add_slide(blank_layout)
add_rect(sl1, 0, 0, 13.33, 7.5, NAVY)
# diagonal accent shape
add_rect(sl1, 8.5, 0, 4.83, 7.5, RGBColor(0x26, 0x35, 0x78))
add_rect(sl1, 0, 5.8, 13.33, 1.7, RGBColor(0x16, 0x1C, 0x4A))

# School + course tag
add_text(sl1, "逢甲大學  ·  資訊工程學系  ·  專題研究", 0.8, 0.6, 10, 0.5,
         font_size=13, color=ICE, align=PP_ALIGN.LEFT)

# Main title
add_text(sl1, "K線型態辨識系統", 0.8, 1.5, 10, 1.1,
         font_size=52, color=WHITE, bold=True, align=PP_ALIGN.LEFT)
add_text(sl1, "進度報告", 0.8, 2.55, 10, 0.7,
         font_size=34, color=ICE, bold=False, align=PP_ALIGN.LEFT)

# Subtitle accent line
add_rect(sl1, 0.8, 3.35, 1.5, 0.06, ACCENT)

# Description
add_text(sl1, "基於 Autoencoder + One-Class SVM 的台股K線型態分析與漲跌預測",
         0.8, 3.55, 9.5, 0.5, font_size=15, color=ICE, align=PP_ALIGN.LEFT)

# Team info at bottom
add_text(sl1, "專題生：郭家齊  ·  呂偉湘", 0.8, 5.9, 8, 0.45,
         font_size=14, color=WHITE, align=PP_ALIGN.LEFT)
add_text(sl1, "指導老師：李俊宏 教授", 0.8, 6.35, 8, 0.45,
         font_size=14, color=ICE, align=PP_ALIGN.LEFT)
add_text(sl1, "2026.04", 11.5, 6.35, 1.6, 0.45,
         font_size=14, color=GRAY_M, align=PP_ALIGN.RIGHT)

nav_bar(sl1, 1)

# ─────────────────────────────────────────────────────────────────────────────
# Slide 2: 本次進度目標
# ─────────────────────────────────────────────────────────────────────────────
sl2 = prs.slides.add_slide(blank_layout)
add_rect(sl2, 0, 0, 13.33, 7.5, GRAY_L)
add_rect(sl2, 0, 0, 13.33, 1.3, NAVY)
add_text(sl2, "本次進度目標", 0.6, 0.3, 10, 0.7,
         font_size=30, color=WHITE, bold=True)
add_text(sl2, "Current Progress", 0.6, 0.9, 10, 0.35,
         font_size=13, color=ICE)

goals = [
    ("1", "確定訓練出的K線型態是有效的",
     "驗證 Autoencoder 萃取的潛在空間特徵可區分看漲／看跌型態"),
    ("2", "測試每個型態的預測準確度",
     "引入四大評估指標：Precision、Recall、Accuracy、F1-Score"),
    ("3", "改善K線相似度比對方式",
     "採用歐氏距離（Euclidean Distance）取代餘弦相似度，提升匹配精準度"),
    ("4", "完成Grid Search最佳參數搜尋",
     "跨 54 組參數組合，找出最佳 HOLD_DAYS、RISE_THRESH、DIST_THRESH、MIN_FREQ"),
]

for i, (num, title, desc) in enumerate(goals):
    gx = 0.5 + (i % 2) * 6.4
    gy = 1.55 + (i // 2) * 2.7
    card(sl2, gx, gy, 6.1, 2.4)
    # number badge
    add_rect(sl2, gx + 0.15, gy + 0.2, 0.5, 0.5, NAVY)
    add_text(sl2, num, gx + 0.15, gy + 0.2, 0.5, 0.5,
             font_size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(sl2, title, gx + 0.75, gy + 0.22, 5.2, 0.5,
             font_size=15, color=NAVY, bold=True)
    add_text(sl2, desc, gx + 0.2, gy + 0.85, 5.7, 1.3,
             font_size=13, color=RGBColor(0x44, 0x44, 0x44), wrap=True)

nav_bar(sl2, 2)

# ─────────────────────────────────────────────────────────────────────────────
# Slide 3: 系統架構
# ─────────────────────────────────────────────────────────────────────────────
sl3 = prs.slides.add_slide(blank_layout)
add_rect(sl3, 0, 0, 13.33, 7.5, GRAY_L)
add_rect(sl3, 0, 0, 13.33, 1.3, NAVY)
add_text(sl3, "系統架構", 0.6, 0.3, 10, 0.7, font_size=30, color=WHITE, bold=True)
add_text(sl3, "System Architecture", 0.6, 0.9, 10, 0.35, font_size=13, color=ICE)

steps = [
    ("①  資料收集", "yfinance 下載\nTW50 台灣50\n50檔股票\n2015–2025", NAVY),
    ("②  特徵萃取", "1日7特徵\n2日10特徵\n3日12特徵\n（K線形狀）", ACCENT),
    ("③  Autoencoder", "降維至\n4維潛在空間\nPyTorch MSELoss\n60 Epochs", RGBColor(0x2E, 0x86, 0xAB)),
    ("④  One-Class SVM", "RBF核\nnu=0.1\n看漲／看跌\n種子訓練", RGBColor(0x4C, 0xAF, 0x50)),
    ("⑤  評估驗證", "Precision\nRecall\nAccuracy\nF1-Score", RED),
]

arrow_color = RGBColor(0x99, 0xAA, 0xBB)
box_w = 2.1
box_h = 3.5
start_x = 0.45

for i, (title, body, col) in enumerate(steps):
    bx = start_x + i * (box_w + 0.3)
    # card body
    add_rect(sl3, bx, 1.55, box_w, box_h, WHITE)
    # top color bar
    add_rect(sl3, bx, 1.55, box_w, 0.55, col)
    add_text(sl3, title, bx + 0.05, 1.6, box_w - 0.1, 0.45,
             font_size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    # body text
    for j, line in enumerate(body.split("\n")):
        add_text(sl3, line, bx + 0.1, 2.25 + j * 0.6, box_w - 0.2, 0.55,
                 font_size=13, color=RGBColor(0x22, 0x22, 0x22), align=PP_ALIGN.CENTER)
    # arrow (except last)
    if i < len(steps) - 1:
        ax = bx + box_w + 0.03
        add_rect(sl3, ax, 3.1, 0.22, 0.05, arrow_color)
        add_text(sl3, "▶", ax - 0.06, 3.0, 0.35, 0.3, font_size=14, color=arrow_color, align=PP_ALIGN.CENTER)

# data note
add_text(sl3, "訓練集：2016–2024（108,215 筆）  ·  測試集：2025",
         2.0, 5.35, 9.33, 0.4, font_size=13, color=GRAY_M, align=PP_ALIGN.CENTER)

nav_bar(sl3, 3)

# ─────────────────────────────────────────────────────────────────────────────
# Slide 4: 評估指標說明
# ─────────────────────────────────────────────────────────────────────────────
sl4 = prs.slides.add_slide(blank_layout)
add_rect(sl4, 0, 0, 13.33, 7.5, GRAY_L)
add_rect(sl4, 0, 0, 13.33, 1.3, NAVY)
add_text(sl4, "四大評估指標", 0.6, 0.3, 10, 0.7, font_size=30, color=WHITE, bold=True)
add_text(sl4, "Evaluation Metrics", 0.6, 0.9, 10, 0.35, font_size=13, color=ICE)

metrics = [
    ("Precision\n精確率", "TP ÷ (TP + FP)",
     "被OCSVM預測「看漲」\n且真實報酬 > 2%\n的比例", ACCENT),
    ("Recall\n召回率", "TP ÷ (TP + FN)",
     "真實看漲K線中\n被正確識別出來\n的比例", RGBColor(0x2E, 0x86, 0xAB)),
    ("Accuracy\n準確率", "(TP+TN) ÷ N",
     "所有預測中\n正確預測\n的比例", RGBColor(0x4C, 0xAF, 0x50)),
    ("F1-Score\nF1分數", "2 × P × R ÷ (P+R)",
     "Precision 與 Recall\n的調和平均數\n（平衡指標）", RED),
]

for i, (name, formula, desc, col) in enumerate(metrics):
    mx = 0.45 + i * 3.2
    card(sl4, mx, 1.55, 3.0, 4.5)
    # top color
    add_rect(sl4, mx, 1.55, 3.0, 0.65, col)
    # name
    for j, ln in enumerate(name.split("\n")):
        add_text(sl4, ln, mx + 0.05, 1.6 + j * 0.35, 2.9, 0.38,
                 font_size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    # formula box
    add_rect(sl4, mx + 0.15, 2.35, 2.7, 0.55, ICE)
    add_text(sl4, formula, mx + 0.15, 2.35, 2.7, 0.55,
             font_size=13, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    # description
    for j, ln in enumerate(desc.split("\n")):
        add_text(sl4, ln, mx + 0.1, 3.05 + j * 0.48, 2.8, 0.45,
                 font_size=13, color=RGBColor(0x33, 0x33, 0x33), align=PP_ALIGN.CENTER)

# confusion matrix legend
add_text(sl4, "混淆矩陣對照：  TP = 真陽性　TN = 真陰性　FP = 假陽性　FN = 假陰性",
         1.0, 6.35, 11.33, 0.45, font_size=12, color=GRAY_M, align=PP_ALIGN.CENTER)

nav_bar(sl4, 4)

# ─────────────────────────────────────────────────────────────────────────────
# Slide 5: 參數調整對比
# ─────────────────────────────────────────────────────────────────────────────
sl5 = prs.slides.add_slide(blank_layout)
add_rect(sl5, 0, 0, 13.33, 7.5, GRAY_L)
add_rect(sl5, 0, 0, 13.33, 1.3, NAVY)
add_text(sl5, "參數調整對比", 0.6, 0.3, 10, 0.7, font_size=30, color=WHITE, bold=True)
add_text(sl5, "Parameter Comparison", 0.6, 0.9, 10, 0.35, font_size=13, color=ICE)

params = [
    ("HOLD_DAYS",   "持有天數",     "3 天",   "5 天",   "5日特徵視窗一致"),
    ("RISE_THRESH", "看漲種子門檻",  "3.00%",  "5.00%",  "篩選更強的看漲型態"),
    ("FALL_THRESH", "看跌種子門檻",  "-3.00%", "-5.00%", "篩選更強的看跌型態"),
    ("DIST_THRESH", "歐氏距離門檻",  "2.5",    "3.0",    "放寬相似度範圍"),
    ("MIN_FREQ",    "最低出現頻率",  "3",      "5",      "排除低頻雜訊型態"),
    ("OCSVM_NU",    "邊界寬鬆程度",  "0.3",    "0.1",    "收緊異常偵測邊界"),
]

# Table header
hx, hy = 0.8, 1.5
col_w = [2.5, 2.3, 1.7, 1.7, 3.8]
headers = ["參數名稱", "中文名稱", "原始值", "目前值", "調整原因"]
add_rect(sl5, hx, hy, sum(col_w) + 0.3, 0.5, NAVY)
cx = hx + 0.1
for j, (hd, cw) in enumerate(zip(headers, col_w)):
    add_text(sl5, hd, cx, hy + 0.05, cw, 0.4,
             font_size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    cx += cw

for i, (param, cname, orig, curr, reason) in enumerate(params):
    ry = hy + 0.5 + i * 0.72
    bg = WHITE if i % 2 == 0 else RGBColor(0xF0, 0xF4, 0xFB)
    add_rect(sl5, hx, ry, sum(col_w) + 0.3, 0.7, bg)
    vals = [param, cname, orig, curr, reason]
    cx = hx + 0.1
    for j, (val, cw) in enumerate(zip(vals, col_w)):
        fc = NAVY if j == 0 else (RED if j == 2 else (GREEN if j == 3 else RGBColor(0x22, 0x22, 0x22)))
        add_text(sl5, val, cx, ry + 0.08, cw, 0.55,
                 font_size=12, color=fc, bold=(j == 0), align=PP_ALIGN.CENTER)
        cx += cw

nav_bar(sl5, 5)

# ─────────────────────────────────────────────────────────────────────────────
# Slide 6: Grid Search 結果
# ─────────────────────────────────────────────────────────────────────────────
sl6 = prs.slides.add_slide(blank_layout)
add_rect(sl6, 0, 0, 13.33, 7.5, GRAY_L)
add_rect(sl6, 0, 0, 13.33, 1.3, NAVY)
add_text(sl6, "Grid Search 最佳參數結果", 0.6, 0.3, 10, 0.7, font_size=30, color=WHITE, bold=True)
add_text(sl6, "54 組參數組合全面測試", 0.6, 0.9, 10, 0.35, font_size=13, color=ICE)

# left: search space
card(sl6, 0.5, 1.5, 5.6, 5.3)
add_text(sl6, "搜尋空間", 0.65, 1.65, 5.0, 0.45,
         font_size=16, color=NAVY, bold=True)
add_rect(sl6, 0.65, 2.1, 5.2, 0.04, ICE)

search_params = [
    ("HOLD_DAYS",   "3, 5, 10 天"),
    ("RISE_THRESH", "2%, 3%, 5%"),
    ("DIST_THRESH", "2.0, 2.5, 3.0"),
    ("MIN_FREQ",    "3, 5 次"),
]
for i, (p, v) in enumerate(search_params):
    add_rect(sl6, 0.75, 2.25 + i * 0.9, 2.2, 0.65, ICE)
    add_text(sl6, p, 0.8, 2.3 + i * 0.9, 2.1, 0.55,
             font_size=12, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(sl6, v, 3.05, 2.3 + i * 0.9, 2.9, 0.55,
             font_size=13, color=RGBColor(0x22, 0x22, 0x22))
add_text(sl6, "總計：3 × 3 × 3 × 2 = 54 組", 0.75, 6.05, 5.0, 0.45,
         font_size=13, color=GRAY_M, italic=True)

# right: best result
card(sl6, 6.5, 1.5, 6.3, 5.3)
add_text(sl6, "最佳組合", 6.65, 1.65, 5.8, 0.45,
         font_size=16, color=NAVY, bold=True)
add_rect(sl6, 6.65, 2.1, 5.9, 0.04, ICE)

best = [
    ("HOLD_DAYS",   "5 天"),
    ("RISE_THRESH", "5.0%"),
    ("DIST_THRESH", "3.0"),
    ("MIN_FREQ",    "5 次"),
]
for i, (p, v) in enumerate(best):
    add_rect(sl6, 6.7, 2.25 + i * 0.75, 2.0, 0.6, ICE)
    add_text(sl6, p, 6.75, 2.3 + i * 0.75, 1.9, 0.5,
             font_size=12, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(sl6, v, 8.85, 2.3 + i * 0.75, 4.0, 0.5,
             font_size=15, color=GREEN, bold=True)

# result highlight
add_rect(sl6, 6.7, 5.45, 5.8, 1.1, RGBColor(0xE8, 0xF5, 0xE9))
add_rect(sl6, 6.7, 5.45, 0.08, 1.1, GREEN)
add_text(sl6, "看漲 Precision", 6.9, 5.55, 2.8, 0.45, font_size=13, color=GRAY_M)
add_text(sl6, "33.40%", 9.8, 5.48, 2.5, 0.6, font_size=28, color=GREEN, bold=True)
add_text(sl6, "（最高 Precision 組合）", 6.9, 5.98, 5.5, 0.4, font_size=11, color=GRAY_M)

nav_bar(sl6, 6)

# ─────────────────────────────────────────────────────────────────────────────
# Slide 7: NU 收緊測試
# ─────────────────────────────────────────────────────────────────────────────
sl7 = prs.slides.add_slide(blank_layout)
add_rect(sl7, 0, 0, 13.33, 7.5, GRAY_L)
add_rect(sl7, 0, 0, 13.33, 1.3, NAVY)
add_text(sl7, "OCSVM NU 收緊測試", 0.6, 0.3, 10, 0.7, font_size=30, color=WHITE, bold=True)
add_text(sl7, "逐步收緊 nu 觀察指標變化（nu = 0.1 → 0.05 → 0.01）", 0.6, 0.9, 10, 0.35, font_size=13, color=ICE)

# bull table
add_text(sl7, "看漲結果", 0.6, 1.45, 5.8, 0.4, font_size=15, color=NAVY, bold=True)
headers7 = ["nu", "預測正例%", "Precision", "Recall", "F1"]
col_w7 = [1.0, 1.5, 1.5, 1.5, 1.5]
add_rect(sl7, 0.55, 1.85, sum(col_w7) + 0.15, 0.45, NAVY)
cx = 0.65
for h, cw in zip(headers7, col_w7):
    add_text(sl7, h, cx, 1.9, cw, 0.35, font_size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    cx += cw

bull_data = [
    ("0.1",  "92.6%", "25.98%", "92.41%", "40.55%"),
    ("0.05", "97.3%", "25.86%", "96.69%", "40.81%"),
    ("0.01", "99.7%", "25.97%", "99.46%", "41.18%"),
]
for i, row in enumerate(bull_data):
    ry = 2.3 + i * 0.65
    bg = WHITE if i % 2 == 0 else RGBColor(0xF0, 0xF4, 0xFB)
    add_rect(sl7, 0.55, ry, sum(col_w7) + 0.15, 0.63, bg)
    cx = 0.65
    for j, (val, cw) in enumerate(zip(row, col_w7)):
        fc = NAVY if j == 0 else (GREEN if j == 2 else RGBColor(0x22, 0x22, 0x22))
        add_text(sl7, val, cx, ry + 0.08, cw, 0.45,
                 font_size=13, color=fc, align=PP_ALIGN.CENTER)
        cx += cw

# bear table
add_text(sl7, "看跌結果", 7.0, 1.45, 5.8, 0.4, font_size=15, color=NAVY, bold=True)
add_rect(sl7, 6.95, 1.85, sum(col_w7) + 0.15, 0.45, NAVY)
cx = 7.05
for h, cw in zip(headers7, col_w7):
    add_text(sl7, h, cx, 1.9, cw, 0.35, font_size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    cx += cw

bear_data = [
    ("0.1",  "93.9%", "21.68%", "92.81%", "35.15%"),
    ("0.05", "97.5%", "21.75%", "96.75%", "35.52%"),
    ("0.01", "99.4%", "21.91%", "99.31%", "35.89%"),
]
for i, row in enumerate(bear_data):
    ry = 2.3 + i * 0.65
    bg = WHITE if i % 2 == 0 else RGBColor(0xF0, 0xF4, 0xFB)
    add_rect(sl7, 6.95, ry, sum(col_w7) + 0.15, 0.63, bg)
    cx = 7.05
    for j, (val, cw) in enumerate(zip(row, col_w7)):
        fc = NAVY if j == 0 else (RED if j == 2 else RGBColor(0x22, 0x22, 0x22))
        add_text(sl7, val, cx, ry + 0.08, cw, 0.45,
                 font_size=13, color=fc, align=PP_ALIGN.CENTER)
        cx += cw

# conclusion box
add_rect(sl7, 0.55, 4.4, 12.23, 1.8, RGBColor(0xFF, 0xF3, 0xCD))
add_rect(sl7, 0.55, 4.4, 0.1, 1.8, GOLD)
add_text(sl7, "結論", 0.8, 4.5, 1.0, 0.4, font_size=14, color=NAVY, bold=True)
conclusions = [
    "• nu 從 0.1 收緊至 0.01，Precision 幾乎沒有變化（約 25%～26%）",
    "• nu 越小反而造成 OCSVM 預測「幾乎所有K線為正例」（99.7%），問題非邊界寬緊",
    "• 根本原因：純K線形狀特徵在潛在空間中，看漲／看跌型態高度重疊 → 需加入技術指標",
]
for j, c in enumerate(conclusions):
    add_text(sl7, c, 0.8, 4.95 + j * 0.42, 12.0, 0.4,
             font_size=13, color=RGBColor(0x55, 0x44, 0x00))

nav_bar(sl7, 7)

# ─────────────────────────────────────────────────────────────────────────────
# Slide 8: 混淆矩陣結果分析
# ─────────────────────────────────────────────────────────────────────────────
sl8 = prs.slides.add_slide(blank_layout)
add_rect(sl8, 0, 0, 13.33, 7.5, GRAY_L)
add_rect(sl8, 0, 0, 13.33, 1.3, NAVY)
add_text(sl8, "混淆矩陣結果分析", 0.6, 0.3, 10, 0.7, font_size=30, color=WHITE, bold=True)
add_text(sl8, "1日 7特徵 | HOLD_DAYS=5, RISE_THRESH=5%, OCSVM_NU=0.1", 0.6, 0.9, 10, 0.35, font_size=13, color=ICE)

# Two big result panels
panels = [
    ("看漲（Bull）", "26,019", "2,137", "74,141", "5,918",
     "25.98%", "92.41%", "29.51%", "40.55%", ACCENT, GREEN),
    ("看跌（Bear）", "22,023", "1,706", "79,539", "4,947",
     "21.68%", "92.81%", "24.92%", "35.15%", RED, RGBColor(0xE5, 0x88, 0x00)),
]

for i, (label, tp, fn, fp, tn, prec, rec, acc, f1, col1, col2) in enumerate(panels):
    px = 0.5 + i * 6.4
    card(sl8, px, 1.5, 6.1, 5.5)
    add_rect(sl8, px, 1.5, 6.1, 0.5, col1)
    add_text(sl8, label, px + 0.1, 1.55, 5.9, 0.4, font_size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    # confusion matrix 2x2
    add_text(sl8, "預測正", px + 1.65, 2.12, 1.5, 0.35, font_size=11, color=GRAY_M, align=PP_ALIGN.CENTER)
    add_text(sl8, "預測負", px + 3.35, 2.12, 1.5, 0.35, font_size=11, color=GRAY_M, align=PP_ALIGN.CENTER)
    add_text(sl8, "實際正", px + 0.1, 2.55, 1.3, 0.9, font_size=11, color=GRAY_M, align=PP_ALIGN.CENTER)
    add_text(sl8, "實際負", px + 0.1, 3.45, 1.3, 0.9, font_size=11, color=GRAY_M, align=PP_ALIGN.CENTER)

    # TP
    add_rect(sl8, px + 1.5, 2.5, 1.7, 0.9, RGBColor(0xD5, 0xF5, 0xE3))
    add_text(sl8, f"TP\n{tp}", px + 1.5, 2.52, 1.7, 0.86, font_size=14, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    # FN
    add_rect(sl8, px + 3.3, 2.5, 1.7, 0.9, RGBColor(0xFD, 0xED, 0xEC))
    add_text(sl8, f"FN\n{fn}", px + 3.3, 2.52, 1.7, 0.86, font_size=14, color=RED, bold=True, align=PP_ALIGN.CENTER)
    # FP
    add_rect(sl8, px + 1.5, 3.45, 1.7, 0.9, RGBColor(0xFD, 0xED, 0xEC))
    add_text(sl8, f"FP\n{fp}", px + 1.5, 3.47, 1.7, 0.86, font_size=14, color=RED, bold=True, align=PP_ALIGN.CENTER)
    # TN
    add_rect(sl8, px + 3.3, 3.45, 1.7, 0.9, RGBColor(0xD5, 0xF5, 0xE3))
    add_text(sl8, f"TN\n{tn}", px + 3.3, 3.47, 1.7, 0.86, font_size=14, color=GREEN, bold=True, align=PP_ALIGN.CENTER)

    # metrics row
    metric_data = [("Precision", prec), ("Recall", rec), ("Accuracy", acc), ("F1", f1)]
    for j, (mn, mv) in enumerate(metric_data):
        mx2 = px + 0.15 + j * 1.45
        add_rect(sl8, mx2, 4.55, 1.35, 1.2, ICE if j % 2 == 0 else RGBColor(0xE8, 0xEF, 0xFF))
        add_text(sl8, mn, mx2 + 0.05, 4.6, 1.25, 0.4, font_size=11, color=NAVY, align=PP_ALIGN.CENTER)
        add_text(sl8, mv, mx2 + 0.05, 5.0, 1.25, 0.6, font_size=18, color=col2, bold=True, align=PP_ALIGN.CENTER)

nav_bar(sl8, 8)

# ─────────────────────────────────────────────────────────────────────────────
# Slide 9: 結論與下一步
# ─────────────────────────────────────────────────────────────────────────────
sl9 = prs.slides.add_slide(blank_layout)
add_rect(sl9, 0, 0, 13.33, 7.5, NAVY)
add_rect(sl9, 0, 0, 4.5, 7.5, RGBColor(0x12, 0x1A, 0x44))
add_text(sl9, "結論\n與\n下一步", 0.4, 1.5, 3.6, 3.0,
         font_size=32, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
add_text(sl9, "Conclusion &\nNext Steps", 0.4, 4.5, 3.6, 1.5,
         font_size=16, color=ICE, align=PP_ALIGN.CENTER)

# right panel
add_text(sl9, "本次研究成果", 5.0, 0.6, 7.8, 0.5, font_size=18, color=ICE, bold=True)
add_rect(sl9, 5.0, 1.1, 7.8, 0.04, RGBColor(0x3A, 0x4A, 0x88))

results_txt = [
    ("✓", "完成 54 組 Grid Search，最佳 Precision：33.40%（看漲）"),
    ("✓", "OCSVM NU 收緊測試確認：特徵是瓶頸，非邊界問題"),
    ("✓", "引入四大評估指標 + 混淆矩陣視覺化"),
    ("✓", "所有K線相似條件改為「每根都必須」而非平均值"),
]
for i, (mark, txt) in enumerate(results_txt):
    add_rect(sl9, 5.0, 1.3 + i * 0.65, 0.4, 0.45, ACCENT)
    add_text(sl9, mark, 5.0, 1.3 + i * 0.65, 0.4, 0.45,
             font_size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(sl9, txt, 5.5, 1.33 + i * 0.65, 7.2, 0.4,
             font_size=13, color=WHITE)

add_text(sl9, "接下來規劃", 5.0, 4.1, 7.8, 0.5, font_size=18, color=ICE, bold=True)
add_rect(sl9, 5.0, 4.6, 7.8, 0.04, RGBColor(0x3A, 0x4A, 0x88))

next_txt = [
    ("→", "加入技術指標：RSI、MA位置、布林通道（提升 Precision）"),
    ("→", "擴充特徵工程：量價關係、市場強弱度"),
    ("→", "實測 2025 年資料，驗證泛化能力"),
]
for i, (mark, txt) in enumerate(next_txt):
    add_rect(sl9, 5.0, 4.8 + i * 0.65, 0.4, 0.45, GOLD)
    add_text(sl9, mark, 5.0, 4.8 + i * 0.65, 0.4, 0.45,
             font_size=14, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(sl9, txt, 5.5, 4.83 + i * 0.65, 7.2, 0.4,
             font_size=13, color=WHITE)

nav_bar(sl9, 9)

# ── Save ──────────────────────────────────────────────────────────────────────
out = r"C:\Users\alex\Desktop\專題進度\0401\kline_progress_report.pptx"
prs.save(out)
print(f"Saved: {out}")
