"""Generate progress report PPT for Chart GCN project."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

WHITE = RGBColor(255, 255, 255)
BLACK = RGBColor(0, 0, 0)
DARK_GRAY = RGBColor(50, 50, 50)
LIGHT_GRAY = RGBColor(240, 240, 240)
ACCENT = RGBColor(0, 120, 212)


def add_black_bar(slide, left=Inches(0.3), top=Inches(0.8), width=Inches(0.08), height=Inches(1.2)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLACK
    shape.line.fill.background()


def add_bottom_bar(slide):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(12.8), Inches(6.5), Inches(0.15), Inches(1.0)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = BLACK
    shape.line.fill.background()


def set_slide_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


# ===================== Slide 1: Title =====================
slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_slide_bg(slide1)

title_box = slide1.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11), Inches(2.5))
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "目前進度及接下來規劃"
p.font.size = Pt(54)
p.font.bold = True
p.font.color.rgb = BLACK
p.alignment = PP_ALIGN.CENTER

info_box = slide1.shapes.add_textbox(Inches(0.8), Inches(5.5), Inches(5), Inches(1.2))
tf2 = info_box.text_frame
tf2.word_wrap = True
p1 = tf2.paragraphs[0]
p1.text = "專題生:郭家齊 呂偉湘"
p1.font.size = Pt(18)
p1.font.color.rgb = DARK_GRAY
p2 = tf2.add_paragraph()
p2.text = "指導老師:李俊宏"
p2.font.size = Pt(18)
p2.font.color.rgb = DARK_GRAY

add_bottom_bar(slide1)

# ===================== Slide 2: This Week Progress =====================
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide2)
add_black_bar(slide2, top=Inches(0.5), height=Inches(1.0))

header2 = slide2.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(6), Inches(1.0))
tf = header2.text_frame
p = tf.paragraphs[0]
p.text = "這禮拜進度"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = BLACK

# Left side: pipeline description
left_box = slide2.shapes.add_textbox(Inches(0.6), Inches(1.8), Inches(6.5), Inches(5.0))
tf = left_box.text_frame
tf.word_wrap = True

items = [
    ("Chart GCN 模型實作", "對齊論文 Li et al., KBS 2022"),
    ("資料來源", "台股前 100 大 (實際取得 95 檔)，2018-2024"),
    ("Pipeline", "PIP 提取 → Visibility Graph → 子圖選取 → CNN + Self-Attention"),
    ("特徵工程", "9 維技術指標 (SMA, WMA, Momentum, MACD, Williams %R, CCI, K%, D%, RSI)"),
    ("模型改進", "加入 BatchNorm、Residual Attention、macro F1 選模"),
    ("Grid Search", "兩階段搜尋 32 組超參數，最佳: window=130, m_pips=80, N=15, g=5"),
    ("訓練設定", "window=130, m_pips=80, N=15, g=5, epochs=30"),
    ("回測驗證", "2024 年測試期，95 檔股票交易模擬"),
]

for i, (title, desc) in enumerate(items):
    if i > 0:
        p = tf.add_paragraph()
    else:
        p = tf.paragraphs[0]
    p.space_before = Pt(6)
    run = p.add_run()
    run.text = f"{title}: "
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = BLACK
    run2 = p.add_run()
    run2.text = desc
    run2.font.size = Pt(16)
    run2.font.color.rgb = DARK_GRAY

# Right side: pipeline flow chart (text-based)
right_box = slide2.shapes.add_textbox(Inches(7.5), Inches(1.8), Inches(5.5), Inches(5.5))
tf = right_box.text_frame
tf.word_wrap = True

pipeline = """Raw Data 台股 OHLCV
    ↓
PIP Algorithm 關鍵轉折點
    ↓
Natural Visibility Graph
    ↓
Subgraph Extraction 子圖選取
    ↓
Feature Matrix (N×g×F)
    ↓
CNN Layers (Conv1→Conv2→Conv3)
    ↓
FC Layers + Self-Attention
    ↓
Prediction (漲/跌)
    ↓
Backtest 回測驗證"""

p = tf.paragraphs[0]
p.text = pipeline
p.font.size = Pt(16)
p.font.color.rgb = BLACK
p.font.name = "Consolas"

add_black_bar(slide2, left=Inches(7.2), top=Inches(1.8), height=Inches(5.0))

# ===================== Slide 3: Results =====================
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide3)
add_black_bar(slide3, top=Inches(0.5), height=Inches(1.0))

header3 = slide3.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(6), Inches(1.0))
tf = header3.text_frame
p = tf.paragraphs[0]
p.text = "實驗結果"
p.font.size = Pt(40)
p.font.bold = True

# Key metrics display
metrics_box = slide3.shapes.add_textbox(Inches(1), Inches(1.7), Inches(11), Inches(1.8))
tf = metrics_box.text_frame
tf.word_wrap = True

metrics_lines = [
    "Grid Search 最佳參數: window=130, m_pips=80, N=15, g=5",
    "Test Accuracy: 51.49%    Test F1 macro: 0.4887",
    "Test F1 (漲): 0.3729    Test F1 (跌): 0.6045",
    "策略報酬: -5.82%    Benchmark: -3.87%    超額報酬: -1.94%",
]
for i, line in enumerate(metrics_lines):
    if i > 0:
        p = tf.add_paragraph()
    else:
        p = tf.paragraphs[0]
    p.text = line
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = BLACK
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(4)

# Table
rows, cols = 8, 4
tbl_shape = slide3.shapes.add_table(rows, cols, Inches(0.8), Inches(3.8), Inches(11.5), Inches(3.2))
table = tbl_shape.table

table.columns[0].width = Inches(3.0)
table.columns[1].width = Inches(2.5)
table.columns[2].width = Inches(3.0)
table.columns[3].width = Inches(3.0)

headers = ["指標", "數值", "說明", "對應論文"]
data = [
    ["Accuracy", "51.49%", "整體分類準確率", "Table 3"],
    ["Precision (漲)", "45.03%", "預測漲的精確率", "Table 3"],
    ["Recall (漲)", "31.82%", "實際漲被捕捉的比例", "Table 3"],
    ["F1 (漲)", "0.3729", "漲類的 F1 分數", "Table 3"],
    ["F1 (跌)", "0.6045", "跌類的 F1 分數", "Table 3"],
    ["F1 Macro", "0.4887", "兩類 F1 的平均", "Section 4.3"],
    ["超額報酬", "-1.94%", "策略 vs 等權持有", "Section 6.3"],
]

for j, h in enumerate(headers):
    cell = table.cell(0, j)
    cell.text = h
    cell.fill.solid()
    cell.fill.fore_color.rgb = DARK_GRAY
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.size = Pt(14)
        paragraph.font.bold = True
        paragraph.font.color.rgb = WHITE
        paragraph.alignment = PP_ALIGN.CENTER

for i, row_data in enumerate(data):
    for j, val in enumerate(row_data):
        cell = table.cell(i + 1, j)
        cell.text = val
        if i % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_GRAY
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(13)
            paragraph.font.color.rgb = BLACK
            paragraph.alignment = PP_ALIGN.CENTER

# ===================== Slide 4: Backtest Chart =====================
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide4)
add_black_bar(slide4, top=Inches(0.5), height=Inches(1.0))

header4 = slide4.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(6), Inches(1.0))
tf = header4.text_frame
p = tf.paragraphs[0]
p.text = "回測結果"
p.font.size = Pt(40)
p.font.bold = True

img_path = os.path.join(os.path.dirname(__file__), "backtest_gridsearch.png")
if not os.path.exists(img_path):
    img_path = os.path.join(os.path.dirname(__file__), "backtest_100stocks.png")
if not os.path.exists(img_path):
    img_path = os.path.join(os.path.dirname(__file__), "backtest_result.png")

if os.path.exists(img_path):
    slide4.shapes.add_picture(img_path, Inches(1.0), Inches(1.7), Inches(11.0), Inches(5.3))
else:
    note = slide4.shapes.add_textbox(Inches(2), Inches(3), Inches(9), Inches(2))
    tf = note.text_frame
    p = tf.paragraphs[0]
    p.text = "(回測圖片未找到)"
    p.font.size = Pt(24)
    p.alignment = PP_ALIGN.CENTER

# Top performers box
top_box = slide4.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(3.0), Inches(2.5))
tf = top_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "表現最佳個股"
p.font.size = Pt(16)
p.font.bold = True
top_stocks = [
    "2615.TW +25.32%",
    "3653.TW +18.16%",
    "6415.TW +15.11%",
    "1590.TW +14.82%",
]
for s in top_stocks:
    p = tf.add_paragraph()
    p.text = f"  {s}"
    p.font.size = Pt(13)
    p.font.color.rgb = DARK_GRAY

# ===================== Slide 5: Future Plans =====================
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide5)
add_black_bar(slide5, top=Inches(0.5), height=Inches(1.2))

header5 = slide5.shapes.add_textbox(Inches(0.6), Inches(0.5), Inches(8), Inches(1.0))
tf = header5.text_frame
p = tf.paragraphs[0]
p.text = "未來一周進度"
p.font.size = Pt(40)
p.font.bold = True

plans_box = slide5.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11), Inches(5.0))
tf = plans_box.text_frame
tf.word_wrap = True

plans = [
    "改進回測策略 (加入停損/停利機制)",
    "嘗試增加 epochs 或調整 learning rate 提升模型表現",
    "與主分支 (AE + OC-SVM) 整合討論",
]

for i, plan in enumerate(plans):
    if i > 0:
        p = tf.add_paragraph()
    else:
        p = tf.paragraphs[0]
    p.text = f"○  {plan}"
    p.font.size = Pt(22)
    p.font.color.rgb = BLACK
    p.space_before = Pt(12)

add_black_bar(slide5, left=Inches(0.3), top=Inches(1.8), height=Inches(4.5))

# Save
out_path = os.path.join(os.path.dirname(__file__), "progress_report.pptx")
prs.save(out_path)
print(f"PPT saved: {out_path}")
