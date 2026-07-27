# -*- coding: utf-8 -*-
"""產生 725/進度報告.pptx — 白底黑字; 只呈現本週「不同檔數」與「分族群」測試結果."""
import os

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

BLACK = RGBColor(0, 0, 0)
GRAY = RGBColor(0x52, 0x51, 0x4E)
FONT = "微軟正黑體"
HERE = os.path.dirname(os.path.abspath(__file__))

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def add_title(slide, text, size=32):
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.1), Inches(1.0))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.name = FONT; r.font.size = Pt(size); r.font.bold = True
    r.font.color.rgb = BLACK


def add_note(slide, text, top=6.6, size=16, color=GRAY):
    tb = slide.shapes.add_textbox(Inches(0.8), Inches(top), Inches(11.8), Inches(0.7))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = text
    r.font.name = FONT; r.font.size = Pt(size); r.font.color.rgb = color


def add_table(slide, rows, top=1.6, left=0.8, width=11.7, col_widths=None, font_size=15):
    n_r, n_c = len(rows), len(rows[0])
    shp = slide.shapes.add_table(n_r, n_c, Inches(left), Inches(top), Inches(width),
                                 Inches(0.42 * n_r))
    tbl = shp.table
    if col_widths:
        for j, w in enumerate(col_widths):
            tbl.columns[j].width = Inches(w)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF) if i else RGBColor(0xEF, 0xEE, 0xEA)
            p = cell.text_frame.paragraphs[0]
            r = p.add_run(); r.text = str(val)
            r.font.name = FONT; r.font.size = Pt(font_size)
            r.font.bold = (i == 0)
            r.font.color.rgb = BLACK


# 1. 標題
s = prs.slides.add_slide(BLANK)
add_title(s, "進度報告:不同檔數與分族群測試結果", size=40)
tb = s.shapes.add_textbox(Inches(0.6), Inches(3.2), Inches(12), Inches(1.5))
for txt, size in [("2026 / 07 / 27", 24),
                  ("協定:date 批次(無洩漏)· raw 特徵 · w140 / m80 / N10 / g4 · seed 42", 16)]:
    p = tb.text_frame.add_paragraph()
    r = p.add_run(); r.text = txt
    r.font.name = FONT; r.font.size = Pt(size); r.font.color.rgb = GRAY

# 2. 不同檔數 — 曲線圖
s = prs.slides.add_slide(BLANK)
add_title(s, "不同檔數:50 → 500(每次 +50, 共 10 輪)")
img = os.path.join(HERE, "scale_tw500_curve.png")
if os.path.exists(img):
    s.shapes.add_picture(img, Inches(1.35), Inches(1.3), width=Inches(10.6))

# 3. 不同檔數 — 結果表
s = prs.slides.add_slide(BLANK)
add_title(s, "不同檔數:結果")
add_table(s, [
    ["檔數", "訓練樣本", "Acc", "F1 macro", "回測超額", "形態"],
    ["50", "90,250", "54.33%", "38.44%", "-10.6%", "塌縮(全押跌)"],
    ["100", "178,846", "54.81%", "47.38%", "+0.1%", "半塌縮"],
    ["150", "269,096", "55.05%", "39.87%", "-2.1%", "塌縮"],
    ["200", "359,346", "52.24%", "43.28%", "-13.5%", "半塌縮"],
    ["250", "448,919", "55.13%", "37.97%", "-7.3%", "塌縮"],
    ["300", "537,529", "52.02%", "43.62%", "-11.9%", "半塌縮"],
    ["350", "625,973", "55.90%", "36.08%", "-7.3%", "全塌縮"],
    ["400", "714,749", "46.22%", "45.69%", "-11.8%", "反向(偏押漲)"],
    ["450", "796,711", "56.16%", "36.00%", "-7.5%", "全塌縮"],
    ["500", "885,162", "56.09%", "35.93%", "-8.8%", "全塌縮(零漲預測)"],
], top=1.45, col_widths=[1.2, 1.9, 1.7, 1.9, 1.8, 3.2], font_size=14)
add_note(s, "樣本增加 10 倍無改善趨勢;F1 macro 在「全押跌 ~36%」與「均衡亂猜 ~47%」兩態間震盪", top=6.55, size=16, color=BLACK)

# 4. 分族群 — 結果表
s = prs.slides.add_slide(BLANK)
add_title(s, "分族群:結果")
add_note(s, "電子族群盤點:上市電子相關共 455 檔(零組件 104、半導體 96、光電 68、電腦週邊 64、其他 46、通信 45、通路 21、資服 11)",
         top=1.35, size=16, color=BLACK)
add_table(s, [
    ["族群", "檔數", "Acc", "F1 macro", "F1(漲)", "回測超額", "形態"],
    ["全電子", "450", "51.29%", "51.22%", "49.4%", "+0.4%", "均衡 ≈ 亂猜(歷來最高 F1m)"],
    ["半導體", "92", "44.85%", "38.24%", "58.4%", "-11.6%", "反向塌縮(全押漲)"],
    ["電子零組件", "104", "49.33%", "48.91%", "53.6%", "-3.8%", "均衡亂猜"],
    ["金融保險(對照)", "32", "54.16%", "39.09%", "8.8%", "-18.4%", "塌縮(全押跌)"],
], top=2.4, col_widths=[2.4, 1.1, 1.4, 1.6, 1.4, 1.5, 2.3], font_size=15)
add_note(s, "沒有任何一組同時達成「Acc 高於類別先驗 + 預測均衡」;同質性最高的半導體反而最差", top=5.3, size=16, color=BLACK)

out = os.path.join(HERE, "進度報告.pptx")
prs.save(out)
print("saved", out)
